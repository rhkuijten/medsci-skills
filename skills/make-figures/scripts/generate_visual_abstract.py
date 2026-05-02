#!/usr/bin/env python3
"""
Visual Abstract Generator for Medical Research Papers.
Fills journal-specific PPTX templates with study content using python-pptx.

Usage:
    python generate_visual_abstract.py \
      --template european_radiology.pptx \
      --title "Effect of smoking on biopsy outcomes" \
      --hypothesis "What is the association of smoking with PTNB outcomes?" \
      --methods "Retrospective cohort|N=1200 patients|Logistic regression" \
      --finding "Smoking was associated with higher complication rates" \
      --citation "Journal Name (2026) FirstAuthor Last et al; DOI: 10.xxxx/..." \
      --visual figures/fig1_roc_curve.png \
      --badges "N=1200|CT chest|Single-center" \
      --output visual_abstract.pptx

    # With PNG conversion (requires LibreOffice):
    python generate_visual_abstract.py ... --output visual_abstract.pptx --png

Templates are stored in:
    ${CLAUDE_SKILL_DIR}/references/visual_abstract_templates/

If --template is a bare name (no path), the script looks for it in the templates directory.
If not found, falls back to medsci_default.pptx.
"""

import argparse
import os
import re
import subprocess
import sys
from pathlib import Path


# --- Shape matching patterns ---
# Each field maps to a list of substrings to search in shape.text (case-insensitive).
FIELD_PATTERNS = {
    "title": ["articletitle"],
    "hypothesis": ["hypothesis", "question"],
    "methods": ["methodology", "flowchart", "bullet point"],
    "visual": ["visual element", "image/illustration", "illustration/graph", "visualelement"],
    "finding": ["main finding", "relevance statement", "main result"],
    "citation": ["authors names", "doi", "eur radiol (year)", "journal (year)", "articlecitation"],
    "badge_patient": ["patient cohort", "patient"],
    "badge_modality": ["modality", "organ"],
    "badge_center": ["single", "multi-center", "center"],
    "footer": ["footernote"],
}


# --- Central Illustration validation rules (Fuster-Mann 2019, JACC 74:2816) ---
# See references/jacc_central_illustration_principles.md for full rationale.
CI_FORBIDDEN_METHODS_TERMS = (
    "cohort flow",
    "inclusion criteria",
    "exclusion criteria",
    "study design",
    "enrollment",
    "randomized",
    "sample size",
    "consort",
    "prisma",
    "stard",
)


def validate_central_illustration(zones: int, label_words: int,
                                  numerical_points: int, raw_text: str,
                                  allow_overrides: list[str] | None = None) -> tuple[bool, list[str]]:
    """Apply CI mode validation rules. Returns (passes, reasons)."""
    allow = set(allow_overrides or [])
    failures: list[str] = []

    if "zones" not in allow and zones > 3:
        failures.append(
            f"more than 3 visual zones (got {zones}); Fuster-Mann: 'simplicity is superior'"
        )
    if "words" not in allow and label_words > 30:
        failures.append(
            f"label word count {label_words} > 30; Fuster-Mann: 'avoid using too much text'"
        )
    if "numerical" not in allow and numerical_points > 4:
        failures.append(
            f"{numerical_points} numerical highlights > 4; Fuster-Mann: 'avoid secondary messages'"
        )
    if "methods" not in allow:
        lower = (raw_text or "").lower()
        offenders = [t for t in CI_FORBIDDEN_METHODS_TERMS if t in lower]
        if offenders:
            failures.append(
                f"methodology terms detected ({', '.join(offenders)}); CI is not a Visual Abstract"
            )
    return (len(failures) == 0, failures)


def match_field(shape_text: str) -> str | None:
    """Return the field name that matches the shape's text, or None."""
    text_lower = shape_text.lower().strip()
    if not text_lower:
        return None
    for field, patterns in FIELD_PATTERNS.items():
        for pattern in patterns:
            if pattern in text_lower:
                return field
    return None


def replace_shape_text(shape, new_text: str, preserve_formatting: bool = True):
    """Replace all text in a shape while preserving font formatting."""
    tf = shape.text_frame
    if not tf.paragraphs:
        return

    if preserve_formatting and tf.paragraphs[0].runs:
        # Preserve the first run's formatting
        first_run = tf.paragraphs[0].runs[0]
        font_name = first_run.font.name
        font_size = first_run.font.size
        font_bold = first_run.font.bold
        try:
            font_color = first_run.font.color.rgb
        except (AttributeError, TypeError):
            font_color = None

        # Handle multi-line text (split by \n)
        lines = new_text.split("\n")
        # Clear all existing paragraphs except first
        for i in range(len(tf.paragraphs) - 1, 0, -1):
            p = tf.paragraphs[i]._p
            p.getparent().remove(p)

        # Set first paragraph
        para = tf.paragraphs[0]
        para.clear()
        run = para.add_run()
        run.text = lines[0]
        if font_name:
            run.font.name = font_name
        if font_size:
            run.font.size = font_size
        if font_bold is not None:
            run.font.bold = font_bold
        if font_color:
            run.font.color.rgb = font_color

        # Add remaining lines as new paragraphs
        for line in lines[1:]:
            from pptx.oxml.ns import qn
            from copy import deepcopy
            new_p = deepcopy(para._p)
            new_p.clear()
            # Re-add run with text
            from lxml import etree
            r_elem = deepcopy(para._p.findall(qn("a:r"))[0]) if para._p.findall(qn("a:r")) else None
            if r_elem is not None:
                for t in r_elem.findall(qn("a:t")):
                    t.text = line
                new_p.append(r_elem)
            tf._txBody.append(new_p)
    else:
        # Simple replacement without formatting preservation
        tf.paragraphs[0].text = new_text


def insert_image_into_shape(slide, shape, image_path: str):
    """Replace a shape with an image, maintaining the shape's position and size."""
    from pptx.util import Emu
    from PIL import Image

    # Get shape bounds
    left = shape.left
    top = shape.top
    box_w = shape.width
    box_h = shape.height

    # Calculate image dimensions maintaining aspect ratio
    with Image.open(image_path) as img:
        img_w, img_h = img.size
        aspect = img_w / img_h
        box_aspect = box_w / box_h

        if aspect > box_aspect:
            # Image is wider than box — fit to width
            new_w = box_w
            new_h = int(box_w / aspect)
            new_left = left
            new_top = top + (box_h - new_h) // 2  # center vertically
        else:
            # Image is taller — fit to height
            new_h = box_h
            new_w = int(box_h * aspect)
            new_left = left + (box_w - new_w) // 2  # center horizontally
            new_top = top

    # Clear shape text
    if hasattr(shape, "text_frame"):
        for para in shape.text_frame.paragraphs:
            para.clear()

    # Add image on top of shape
    slide.shapes.add_picture(image_path, new_left, new_top, new_w, new_h)


def fill_template(template_path: str, content: dict, output_path: str,
                  visual_path: str | None = None, slide_index: int = 0):
    """Fill a PPTX template with content and save."""
    from pptx import Presentation

    prs = Presentation(template_path)

    if slide_index >= len(prs.slides):
        print(f"Error: Slide index {slide_index} not found in template "
              f"(template has {len(prs.slides)} slides)", file=sys.stderr)
        sys.exit(1)

    slide = prs.slides[slide_index]
    matched_fields = set()

    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue

        field = match_field(shape.text)
        if field is None:
            continue

        if field == "visual" and visual_path and Path(visual_path).exists():
            insert_image_into_shape(slide, shape, visual_path)
            matched_fields.add(field)
        elif field in content and content[field]:
            replace_shape_text(shape, content[field])
            matched_fields.add(field)

    # Report matching results
    all_fields = set(content.keys())
    if visual_path:
        all_fields.add("visual")
    unmatched = all_fields - matched_fields
    if unmatched:
        print(f"Warning: No matching shape found for: {', '.join(unmatched)}",
              file=sys.stderr)

    # Remove unused slides (keep only the target slide)
    # python-pptx slide deletion requires direct XML manipulation
    sldIdLst = prs.slides._sldIdLst
    indices_to_remove = sorted(
        [i for i in range(len(sldIdLst)) if i != slide_index],
        reverse=True  # remove from end to preserve indices
    )
    for i in indices_to_remove:
        sldId = sldIdLst[i]
        rId = sldId.rId
        sldIdLst.remove(sldId)
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass  # relationship already cleaned up

    prs.save(output_path)
    print(f"Visual abstract saved: {output_path}", file=sys.stderr)
    print(output_path)


def convert_to_png(pptx_path: str, dpi: int = 300) -> str | None:
    """Convert PPTX to PNG using LibreOffice CLI. Returns PNG path or None."""
    output_dir = str(Path(pptx_path).parent)
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "png", pptx_path,
             "--outdir", output_dir],
            check=True, capture_output=True, timeout=30
        )
        png_path = str(Path(pptx_path).with_suffix(".png"))
        if Path(png_path).exists():
            print(f"PNG exported: {png_path}", file=sys.stderr)
            return png_path
    except FileNotFoundError:
        print("Warning: LibreOffice not found. PPTX saved but PNG conversion skipped.",
              file=sys.stderr)
        print("Install with: brew install --cask libreoffice", file=sys.stderr)
    except subprocess.TimeoutExpired:
        print("Warning: LibreOffice conversion timed out.", file=sys.stderr)
    except subprocess.CalledProcessError as e:
        print(f"Warning: LibreOffice conversion failed: {e.stderr.decode()[:200]}",
              file=sys.stderr)
    return None


def resolve_template(template_arg: str) -> str:
    """Resolve template name to full path."""
    # If it's already a full path, use it
    if os.path.isabs(template_arg) and Path(template_arg).exists():
        return template_arg

    # Look in the skill's template directory
    skill_dir = os.environ.get("CLAUDE_SKILL_DIR", "")
    if skill_dir:
        templates_dir = Path(skill_dir) / "references" / "visual_abstract_templates"
    else:
        templates_dir = Path(__file__).parent.parent / "references" / "visual_abstract_templates"

    # Try exact match
    candidate = templates_dir / template_arg
    if candidate.exists():
        return str(candidate)

    # Try with .pptx extension
    candidate = templates_dir / f"{template_arg}.pptx"
    if candidate.exists():
        return str(candidate)

    # Fallback to medsci_default
    default = templates_dir / "medsci_default.pptx"
    if default.exists():
        print(f"Template '{template_arg}' not found, using medsci_default.pptx",
              file=sys.stderr)
        return str(default)

    print(f"Error: No template found. Searched in: {templates_dir}", file=sys.stderr)
    sys.exit(1)


def parse_methods(methods_str: str) -> str:
    """Convert pipe-separated methods to bullet-point format."""
    items = [m.strip() for m in methods_str.split("|") if m.strip()]
    if len(items) <= 1:
        return methods_str
    return "\n".join(f"• {item}" for item in items)


def main():
    parser = argparse.ArgumentParser(
        description="Generate visual abstract from PPTX template"
    )
    parser.add_argument("--type", choices=["visual-abstract", "central-illustration"],
                        default="visual-abstract",
                        help="Artifact type: 'visual-abstract' (methods+results) or "
                             "'central-illustration' (single key finding, JACC house style)")
    parser.add_argument("--template", "-t", default=None,
                        help="Template name or path. Default depends on --type: "
                             "visual-abstract→medsci_default, central-illustration→jacc_central_illustration")
    parser.add_argument("--title", help="Article title (required for VA mode)")
    parser.add_argument("--hypothesis", help="Research question or hypothesis (VA mode)")
    parser.add_argument("--methods", help="Methodology (pipe-separated for bullets, VA mode)")
    parser.add_argument("--finding", help="Main finding (<20 words, VA mode)")
    parser.add_argument("--citation", help="Citation line (journal, year, authors, DOI)")
    parser.add_argument("--visual", help="Path to visual element image (PNG/JPG)")
    parser.add_argument("--badges", help="Three pipe-separated badge texts: cohort|modality|center (VA mode)")
    parser.add_argument("--output", "-o", required=True, help="Output PPTX path")
    parser.add_argument("--png", action="store_true",
                        help="Also convert to PNG (requires LibreOffice)")
    parser.add_argument("--slide-index", type=int, default=0,
                        help="Template slide index to use (default: 0)")
    parser.add_argument("--ci-zones", type=int, default=None,
                        help="CI mode: declared visual zone count (for validation)")
    parser.add_argument("--ci-label-words", type=int, default=None,
                        help="CI mode: total label word count (for validation)")
    parser.add_argument("--ci-numerical-points", type=int, default=None,
                        help="CI mode: count of numerical highlights (for validation)")
    parser.add_argument("--ci-raw-text", default="",
                        help="CI mode: raw text content for forbidden-methods-term scan")
    parser.add_argument("--ci-allow", action="append", default=[],
                        choices=["zones", "words", "numerical", "methods"],
                        help="CI mode: override a single rule (repeatable)")
    args = parser.parse_args()

    # Default template per type
    if args.template is None:
        args.template = ("jacc_central_illustration"
                         if args.type == "central-illustration" else "medsci_default")

    # CI mode: run validation BEFORE generating the PPTX
    if args.type == "central-illustration":
        zones = args.ci_zones if args.ci_zones is not None else 1
        words = args.ci_label_words if args.ci_label_words is not None else 0
        nums = args.ci_numerical_points if args.ci_numerical_points is not None else 0
        passes, reasons = validate_central_illustration(
            zones=zones, label_words=words, numerical_points=nums,
            raw_text=args.ci_raw_text, allow_overrides=args.ci_allow,
        )
        if not passes:
            print("Central Illustration validation FAILED:", file=sys.stderr)
            for r in reasons:
                print(f"  - {r}", file=sys.stderr)
            print("\nSee references/jacc_central_illustration_principles.md for guidance.",
                  file=sys.stderr)
            print("Override with --ci-allow {zones|words|numerical|methods} (use sparingly).",
                  file=sys.stderr)
            sys.exit(2)
        # CI mode does not require --title (only --visual + --citation)
        if not args.visual:
            print("CI mode requires --visual (the author content figure).", file=sys.stderr)
            sys.exit(2)
        # Synthesize a placeholder title so downstream code does not fail
        if not args.title:
            args.title = "(central illustration — title carried by manuscript)"

    # Resolve template
    template_path = resolve_template(args.template)

    # Build content dict
    content = {
        "title": args.title,
        "hypothesis": args.hypothesis or "",
        "methods": parse_methods(args.methods) if args.methods else "",
        "finding": args.finding or "",
        "citation": args.citation or "",
    }

    # Parse badges
    if args.badges:
        badges = [b.strip() for b in args.badges.split("|")]
        content["badge_patient"] = badges[0] if len(badges) > 0 else ""
        content["badge_modality"] = badges[1] if len(badges) > 1 else ""
        content["badge_center"] = badges[2] if len(badges) > 2 else ""

    # Validate visual path
    visual_path = None
    if args.visual:
        if not Path(args.visual).exists():
            print(f"Warning: Visual file not found: {args.visual}", file=sys.stderr)
        else:
            visual_path = str(Path(args.visual).resolve())

    # Ensure output directory exists
    Path(args.output).parent.mkdir(parents=True, exist_ok=True)

    # Fill template
    fill_template(template_path, content, args.output, visual_path, args.slide_index)

    # Optional PNG conversion
    if args.png:
        convert_to_png(args.output)


if __name__ == "__main__":
    main()
